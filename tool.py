"""
title: Edit Office Files
author: giofsp
author_url: https://github.com/sergiofspedro
description: Unified tool to read, edit, and create Office files (.xlsx, .xls, .docx, .pptx) preserving original formatting and styles. Detects highlights, bold, italic formatting. Detects legacy .doc and .ppt. Note: Track changes are not supported.
version: 1.1.0
requirements: openpyxl, python-docx, python-pptx, xlrd
"""
import json
import io
import os
import re
import sqlite3
import sys
import traceback
import platform
from copy import copy
from typing import Optional, List, Dict, Any


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
_DB_PATH = os.path.join(
    os.environ.get("OPEN_WEBUI_DATA_DIR", ""),
    "data", "webui.db",
)
if not os.path.isfile(_DB_PATH):
    _DB_PATH = os.path.join(
        os.path.expanduser("~"),
        "AppData", "Roaming", "open-webui", "data", "webui.db",
    )

_UPLOAD_DIR = os.path.join(
    os.environ.get("OPEN_WEBUI_DATA_DIR", ""),
    "data", "uploads",
)
if not os.path.isdir(_UPLOAD_DIR):
    _UPLOAD_DIR = os.path.join(
        os.path.expanduser("~"),
        "AppData", "Roaming", "open-webui", "data", "uploads",
    )

_EXPORT_DIR = os.environ.get("OWUI_EXPORTS_DIR", os.path.join(os.path.expanduser("~"), "open-webui", "exports"))


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _resolve_file_path(file_id: str) -> Optional[str]:
    """Resolve an Open WebUI file UUID to an absolute disk path."""
    try:
        conn = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
        row = conn.execute(
            "SELECT path FROM file WHERE id = ?", (file_id,)
        ).fetchone()
        conn.close()
    except Exception as exc:
        print(f"[office] DB lookup failed for {file_id}: {exc}", file=sys.stderr)
        return None

    if not row or not row[0]:
        # Fallback: try by filename
        try:
            conn = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
            row = conn.execute(
                "SELECT path FROM file WHERE filename LIKE ?",
                (f"%{file_id}%",),
            ).fetchone()
            conn.close()
        except Exception:
            pass

    if not row or not row[0]:
        print(f"[office] No path for file_id {file_id}", file=sys.stderr)
        return None

    path = row[0]
    if os.path.isfile(path):
        return path

    # Fallback: uploads directory
    candidate = os.path.join(_UPLOAD_DIR, os.path.basename(path))
    if os.path.isfile(candidate):
        return candidate

    # Last resort: UUID prefix match in uploads
    prefix = file_id.split("-")[0] if "-" in file_id else file_id[:8]
    if os.path.isdir(_UPLOAD_DIR):
        for name in os.listdir(_UPLOAD_DIR):
            if name.startswith(prefix):
                candidate = os.path.join(_UPLOAD_DIR, name)
                if os.path.isfile(candidate):
                    return candidate

    print(f"[office] File not found on disk: {path}", file=sys.stderr)
    return None


def _read_file_bytes(file_id: str) -> Optional[bytes]:
    """Return raw bytes for an Open WebUI file."""
    path = _resolve_file_path(file_id)
    if path is None:
        return None
    # Path traversal guard: ensure resolved path stays inside allowed directories
    try:
        _abs = os.path.realpath(path)
        _allowed = False
        for _base in (_UPLOAD_DIR, _EXPORT_DIR,
                       os.path.join(os.environ.get("OPEN_WEBUI_DATA_DIR", ""), "data"),
                       os.path.join(os.path.expanduser("~"), "AppData", "Roaming", "open-webui", "data")):
            if _base and os.path.realpath(_base) in _abs:
                _allowed = True
                break
        if not _allowed:
            print(f"[office] Path traversal blocked: {path}", file=sys.stderr)
            return None
    except Exception:
        pass
    try:
        with open(path, "rb") as fh:
            return fh.read()
    except Exception as exc:
        print(f"[office] Read failed for {path}: {exc}", file=sys.stderr)
        return None


def _detect_type(filename: str) -> str:
    """Detect Office file type from filename."""
    lower = filename.lower()
    if lower.endswith(".xlsx"):
        return "xlsx"
    if lower.endswith(".xls"):
        return "xls"
    if lower.endswith(".docx"):
        return "docx"
    if lower.endswith(".doc"):
        return "doc"
    if lower.endswith(".pptx"):
        return "pptx"
    if lower.endswith(".ppt"):
        return "ppt"
    return "unknown"


def _save_file_sync(file_bytes: bytes, filename: str) -> Optional[str]:
    """Save file to exports dir and return HTTP download URL."""
    try:
        os.makedirs(_EXPORT_DIR, exist_ok=True)
        safe = "".join(c for c in filename if c.isalnum() or c in "._- ")
        if not safe:
            safe = "export"
        base_name, ext = os.path.splitext(safe)
        fname = safe
        counter = 2
        while os.path.exists(os.path.join(_EXPORT_DIR, fname)):
            fname = f"{base_name} ({counter}){ext}"
            counter += 1
        fpath = os.path.join(_EXPORT_DIR, fname)
        with open(fpath, "wb") as f:
            f.write(file_bytes)
        return f"http://localhost:9000/{fname}"
    except Exception as e:
        print(f"[office] Export save failed: {e}", file=sys.stderr)
        return None


def _cell_value(cell):
    """Extract a JSON-safe value from an openpyxl cell."""
    import datetime
    v = cell.value
    if v is None:
        return None
    if isinstance(v, (datetime.datetime, datetime.date)):
        return v.isoformat()
    if isinstance(v, (int, float, str, bool)):
        return v
    return str(v)


def _xls_to_xlsx(xls_data: bytes) -> bytes:
    """Convert .xls bytes to .xlsx bytes using xlrd + openpyxl.

    Returns an in-memory .xlsx workbook as bytes.
    """
    import xlrd
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment

    xls_book = xlrd.open_workbook(file_contents=xls_data)
    xlsx_wb = openpyxl.Workbook()
    # Remove the default sheet; we'll add one per xls sheet
    xlsx_wb.remove(xlsx_wb.active)

    for sheet_idx in range(xls_book.nsheets):
        xls_sheet = xls_book.sheet_by_index(sheet_idx)
        ws = xlsx_wb.create_sheet(title=xls_sheet.name[:31])  # Excel 31-char limit

        for rx in range(xls_sheet.nrows):
            for cx in range(xls_sheet.ncols):
                cell = xls_sheet.cell(rx, cx)
                value = cell.value

                # xlrd date handling: if cell type is XL_CELL_DATE, convert
                if cell.ctype == xlrd.XL_CELL_DATE:
                    try:
                        dt_tuple = xls_book.datemode, int(value)
                        import datetime as _dt
                        value = _dt.datetime(*xlrd.xldate_as_tuple(value, xls_book.datemode))
                    except Exception:
                        pass
                elif cell.ctype == xlrd.XL_CELL_BOOLEAN:
                    value = bool(value)

                ws.cell(row=rx + 1, column=cx + 1, value=value)

        # Auto-fit column widths (rough estimate)
        for col_cells in ws.columns:
            max_len = 0
            for cell in col_cells:
                try:
                    max_len = max(max_len, len(str(cell.value or "")))
                except Exception:
                    pass
            letter = openpyxl.utils.get_column_letter(col_cells[0].column)
            ws.column_dimensions[letter].width = min(max_len + 2, 50)

    out = io.BytesIO()
    xlsx_wb.save(out)
    xlsx_wb.close()
    xls_book.release_resources()
    out.seek(0)
    return out.read()


# =========================================================================
class Tools:
    class Valves:
        export_dir: str = ""
        file_server_url: str = ""

    def __init__(self):
        self.valves = self.Valves()
        ed = self.valves.export_dir or os.path.join(
            os.environ.get("OWUI_EXPORTS_DIR", os.environ.get("LOCALAPPDATA", os.path.expanduser("~"))),
            "open-webui", "exports",
        )
        os.makedirs(ed, exist_ok=True)

    # -----------------------------------------------------------------
    # Internal: save and return markdown link
    # -----------------------------------------------------------------
    async def _save_and_link(self, file_bytes: bytes, filename: str) -> tuple:
        """Save file and return (url, filename) or (None, None)."""
        try:
            ed = self.valves.export_dir or os.path.join(
                os.environ.get("OWUI_EXPORTS_DIR", os.environ.get("LOCALAPPDATA", os.path.expanduser("~"))),
                "open-webui", "exports",
            )
            os.makedirs(ed, exist_ok=True)
            safe = "".join(c for c in filename if c.isalnum() or c in "._- ")
            if not safe:
                safe = "export"
            base_name, ext = os.path.splitext(safe)
            fname = safe
            counter = 2
            while os.path.exists(os.path.join(ed, fname)):
                fname = f"{base_name} ({counter}){ext}"
                counter += 1
            fpath = os.path.join(ed, fname)
            with open(fpath, "wb") as f:
                f.write(file_bytes)
            base_url = (self.valves.file_server_url or "http://localhost:9000").rstrip("/")
            url = f"{base_url}/{fname}"
            return (url, filename)
        except Exception as e:
            print(f"[office] Export save failed: {e}", file=sys.stderr)
            try:
                import base64 as b64
                data = b64.b64encode(file_bytes).decode("ascii")
                ext = os.path.splitext(filename)[1].lower()
                mt = {".xlsx":"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",".docx":"application/vnd.openxmlformats-officedocument.wordprocessingml.document",".pptx":"application/vnd.openxmlformats-officedocument.presentationml.presentation"}
                return (f"data:{mt.get(ext,"application/octet-stream")};base64,{data}", filename)
            except:
                return (None, None)

    # -----------------------------------------------------------------
    # READ
    # -----------------------------------------------------------------
    async def read_file(
        self,
        file_id: str,
        max_rows: int = 500,
        __user__=None,
        __request__=None,
    ) -> str:
        """Read any Office file (.xlsx, .xls, .docx, .pptx) and return its contents as structured JSON.

        Auto-detects the file type from the file ID or filename.
        For xlsx/xls: returns sheets with headers and rows.
        For docx: returns paragraphs with styles and tables.
        For pptx: returns slides with shapes and text.
        Legacy .doc and .ppt formats return a helpful error message.

        Args:
            file_id: The Open WebUI file ID (UUID) or filename
            max_rows: Maximum rows to return for xlsx (default 500)
        """
        try:
            file_data = _read_file_bytes(file_id)
            if file_data is None:
                return json.dumps({
                    "error": (
                        f"Could not read file {file_id}. "
                        "Make sure the file was uploaded via the chat."
                    )
                })

            # Detect type from filename in DB
            try:
                conn = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
                row = conn.execute(
                    "SELECT filename FROM file WHERE id = ?", (file_id,)
                ).fetchone()
                conn.close()
                filename = row[0] if row else file_id
            except Exception:
                filename = file_id

            file_type = _detect_type(filename)
            result: Dict[str, Any] = {
                "file_id": file_id,
                "filename": filename,
                "type": file_type,
            }

            if file_type == "xlsx":
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(file_data), data_only=True)
                result["sheets"] = []
                for sn in wb.sheetnames:
                    ws = wb[sn]
                    sheet: Dict[str, Any] = {
                        "name": sn,
                        "headers": [],
                        "rows": [],
                        "total_rows": ws.max_row or 0,
                        "total_cols": ws.max_column or 0,
                    }
                    max_r = min(ws.max_row or 0, max_rows)
                    for ri, row in enumerate(ws.iter_rows(min_row=1, max_row=max_r), 1):
                        rd = [_cell_value(c) for c in row]
                        if ri == 1:
                            sheet["headers"] = [str(v) if v is not None else "" for v in rd]
                        else:
                            sheet["rows"].append(rd)
                    result["sheets"].append(sheet)
                wb.close()

            elif file_type == "xls":
                import xlrd
                xls_book = xlrd.open_workbook(file_contents=file_data)
                result["sheets"] = []
                for sheet_idx in range(xls_book.nsheets):
                    xls_sheet = xls_book.sheet_by_index(sheet_idx)
                    sheet: Dict[str, Any] = {
                        "name": xls_sheet.name,
                        "headers": [],
                        "rows": [],
                        "total_rows": xls_sheet.nrows,
                        "total_cols": xls_sheet.ncols,
                    }
                    max_r = min(xls_sheet.nrows, max_rows)
                    for rx in range(max_r):
                        row_values = []
                        for cx in range(xls_sheet.ncols):
                            cell = xls_sheet.cell(rx, cx)
                            value = cell.value
                            if cell.ctype == xlrd.XL_CELL_DATE:
                                try:
                                    import datetime as _dt
                                    value = _dt.datetime(*xlrd.xldate_as_tuple(value, xls_book.datemode))
                                except Exception:
                                    pass
                            elif cell.ctype == xlrd.XL_CELL_BOOLEAN:
                                value = bool(value)
                            row_values.append(value)
                        if rx == 0:
                            sheet["headers"] = [str(v) if v is not None else "" for v in row_values]
                        else:
                            sheet["rows"].append(row_values)
                    result["sheets"].append(sheet)
                xls_book.release_resources()

            elif file_type == "docx":
                from docx import Document
                from docx.enum.text import WD_COLOR_INDEX
                doc = Document(io.BytesIO(file_data))
                paragraphs = []
                for p in doc.paragraphs:
                    if p.text.strip():
                        style = p.style.name if p.style else "Normal"
                        runs_info = []
                        for run in p.runs:
                            run_data = {"text": run.text}
                            if run.font.highlight_color and run.font.highlight_color != WD_COLOR_INDEX.AUTO:
                                run_data["highlighted"] = True
                                run_data["highlight_color"] = str(run.font.highlight_color)
                            if run.font.bold:
                                run_data["bold"] = True
                            if run.font.italic:
                                run_data["italic"] = True
                            runs_info.append(run_data)
                        paragraphs.append({"style": style, "text": p.text, "runs": runs_info})
                tables = []
                for t in doc.tables:
                    tbl = {"rows": []}
                    for row in t.rows:
                        tbl["rows"].append([cell.text for cell in row.cells])
                    tables.append(tbl)
                result["paragraphs"] = paragraphs
                result["tables"] = tables

            elif file_type == "pptx":
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_data))
                slides = []
                for si, slide in enumerate(prs.slides, 1):
                    sdata: Dict[str, Any] = {"number": si, "shapes": []}
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            sdata["shapes"].append({
                                "type": str(shape.shape_type),
                                "name": shape.name,
                                "text": shape.text[:500],
                            })
                        if shape.has_table:
                            tbl = {"rows": []}
                            for row in shape.table.rows:
                                tbl["rows"].append([cell.text for cell in row.cells])
                            sdata["tables"] = tbl
                    slides.append(sdata)
                result["slides"] = slides

            elif file_type == "doc":
                result["error"] = "Legacy .doc format is not supported. Please convert to .docx first."

            elif file_type == "ppt":
                result["error"] = "Legacy .ppt format is not supported. Please convert to .pptx first."

            else:
                result["error"] = f"Unsupported file type. Detected: {file_type}. Supported: xlsx, xls, docx, pptx"

            return json.dumps(result, indent=2, default=str, ensure_ascii=False)

        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    # -----------------------------------------------------------------
    # ADD CONTENT
    # -----------------------------------------------------------------
    async def add_content(
        self,
        file_id: str,
        content: str,
        output_filename: str = "",
        __user__=None,
        __request__=None,
    ) -> str:
        """Add new content to an Office file while preserving original formatting.

        For spreadsheets (xlsx): content is CSV text with rows to add.
        For documents (docx): content is text to append at the end.
        For presentations (pptx): each line defines a new slide. Use "---" as separator between slides.

        Args:
            file_id: File ID to edit
            content: Content to add (CSV for xlsx, text for docx/pptx)
            output_filename: Optional output filename
        """
        try:
            file_data = _read_file_bytes(file_id)
            if file_data is None:
                return json.dumps({"error": f"Could not read file {file_id}"})

            # Detect type
            try:
                conn = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
                row = conn.execute(
                    "SELECT filename FROM file WHERE id = ?", (file_id,)
                ).fetchone()
                conn.close()
                filename = row[0] if row else file_id
            except Exception:
                filename = file_id

            file_type = _detect_type(filename)
            out = io.BytesIO()
            out_name = output_filename

            if file_type == "xlsx":
                import openpyxl
                from openpyxl.styles import Font, PatternFill, Alignment
                wb = openpyxl.load_workbook(io.BytesIO(file_data))
                ws = wb.active
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.xlsx"

                # Parse CSV content
                import csv as _csv_mod
                reader = _csv_mod.reader(io.StringIO(content))
                parsed_rows = []
                for csv_row in reader:
                    converted = []
                    for v in csv_row:
                        v = v.strip()
                        if v == '':
                            converted.append(None)
                        else:
                            try:
                                converted.append(int(v))
                            except ValueError:
                                try:
                                    converted.append(float(v))
                                except ValueError:
                                    if v.lower() == 'true':
                                        converted.append(True)
                                    elif v.lower() == 'false':
                                        converted.append(False)
                                    else:
                                        converted.append(v)
                    parsed_rows.append(converted)

                if not parsed_rows:
                    return json.dumps({"error": "No rows provided in CSV content"})

                # Get reference styles from last row
                ref = {}
                if ws.max_row and ws.max_row >= 1:
                    for cell in ws[ws.max_row]:
                        if cell.has_style:
                            ref[cell.column] = {
                                "font": copy(cell.font),
                                "fill": copy(cell.fill),
                                "border": copy(cell.border),
                                "alignment": copy(cell.alignment),
                                "number_format": cell.number_format,
                            }

                start = (ws.max_row or 0) + 1
                for i, rd in enumerate(parsed_rows):
                    for j, v in enumerate(rd, 1):
                        cell = ws.cell(row=start + i, column=j)
                        if j in ref:
                            try:
                                cell.font = copy(ref[j]["font"])
                                cell.fill = copy(ref[j]["fill"])
                                cell.border = copy(ref[j]["border"])
                                cell.alignment = copy(ref[j]["alignment"])
                                cell.number_format = ref[j]["number_format"]
                            except Exception:
                                pass
                        cell.value = v

                wb.save(out)
                wb.close()

            elif file_type == "xls":
                # Convert .xls to .xlsx, then apply same add logic
                file_data = _xls_to_xlsx(file_data)
                file_type = "xlsx"
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.xlsx"
                import openpyxl
                from openpyxl.styles import Font, PatternFill, Alignment
                wb = openpyxl.load_workbook(io.BytesIO(file_data))
                ws = wb.active

                # Parse CSV content
                import csv as _csv_mod
                reader = _csv_mod.reader(io.StringIO(content))
                parsed_rows = []
                for csv_row in reader:
                    converted = []
                    for v in csv_row:
                        v = v.strip()
                        if v == '':
                            converted.append(None)
                        else:
                            try:
                                converted.append(int(v))
                            except ValueError:
                                try:
                                    converted.append(float(v))
                                except ValueError:
                                    if v.lower() == 'true':
                                        converted.append(True)
                                    elif v.lower() == 'false':
                                        converted.append(False)
                                    else:
                                        converted.append(v)
                    parsed_rows.append(converted)

                if not parsed_rows:
                    return json.dumps({"error": "No rows provided in CSV content"})

                ref = {}
                if ws.max_row and ws.max_row >= 1:
                    for cell in ws[ws.max_row]:
                        if cell.has_style:
                            ref[cell.column] = {
                                "font": copy(cell.font),
                                "fill": copy(cell.fill),
                                "border": copy(cell.border),
                                "alignment": copy(cell.alignment),
                                "number_format": cell.number_format,
                            }

                start = (ws.max_row or 0) + 1
                for i, rd in enumerate(parsed_rows):
                    for j, v in enumerate(rd, 1):
                        cell = ws.cell(row=start + i, column=j)
                        if j in ref:
                            try:
                                cell.font = copy(ref[j]["font"])
                                cell.fill = copy(ref[j]["fill"])
                                cell.border = copy(ref[j]["border"])
                                cell.alignment = copy(ref[j]["alignment"])
                                cell.number_format = ref[j]["number_format"]
                            except Exception:
                                pass
                        cell.value = v

                wb.save(out)
                wb.close()

            elif file_type == "docx":
                from docx import Document
                doc = Document(io.BytesIO(file_data))
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.docx"

                # Append paragraphs, preserving last paragraph's style
                last_style = "Normal"
                if doc.paragraphs:
                    last_style = doc.paragraphs[-1].style.name if doc.paragraphs[-1].style else "Normal"

                for line in content.split("\n"):
                    doc.add_paragraph(line, style=last_style)

                out = io.BytesIO()
                doc.save(out)

            elif file_type == "pptx":
                from pptx import Presentation
                from pptx.util import Inches
                prs = Presentation(io.BytesIO(file_data))
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.pptx"

                # Split content by "---" into slides
                slide_specs = re.split(r'\n---\n|\r\n---\r\n|\n---\n', content)
                blank_layout = prs.slide_layouts[6]  # Blank layout

                for spec in slide_specs:
                    spec = spec.strip()
                    if not spec:
                        continue
                    lines = spec.split("\n")
                    title = lines[0].strip()
                    body = "\n".join(lines[1:]).strip() if len(lines) > 1 else ""

                    slide = prs.slides.add_slide(blank_layout)
                    # Add title
                    txBox = slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.3), Inches(9), Inches(1)
                    )
                    tf = txBox.text_frame
                    tf.text = title
                    p = tf.paragraphs[0]
                    p.font.size = Inches(0.6)
                    p.font.bold = True

                    # Add body text
                    if body:
                        txBox2 = slide.shapes.add_textbox(
                            Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)
                        )
                        tf2 = txBox2.text_frame
                        tf2.text = body
                        for para in tf2.paragraphs:
                            para.font.size = Inches(0.3)

                out = io.BytesIO()
                prs.save(out)

            else:
                return json.dumps({"error": f"Unsupported type: {file_type}"})

            out.seek(0)
            url, name = await self._save_and_link(out.read(), out_name)
            if url:
                return f"[{name}]({url})\n\nAdded content to {file_type.upper()} file, preserving original formatting."
            return json.dumps({"error": "Could not save file"})

        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    # -----------------------------------------------------------------
    # REPLACE TEXT
    # -----------------------------------------------------------------
    async def replace_text(
        self,
        file_id: str,
        find_text: str,
        replace_with: str,
        output_filename: str = "",
        __user__=None,
        __request__=None,
    ) -> str:
        """Find and replace text in any Office file while preserving original formatting.

        Works on cell values in xlsx, paragraph text in docx, and shape text in pptx.

        Args:
            file_id: File ID to edit
            find_text: Text to find
            replace_with: Text to replace with
            output_filename: Optional output filename
        """
        try:
            file_data = _read_file_bytes(file_id)
            if file_data is None:
                return json.dumps({"error": f"Could not read file {file_id}"})

            # Detect type
            try:
                conn = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
                row = conn.execute(
                    "SELECT filename FROM file WHERE id = ?", (file_id,)
                ).fetchone()
                conn.close()
                filename = row[0] if row else file_id
            except Exception:
                filename = file_id

            file_type = _detect_type(filename)
            out = io.BytesIO()
            out_name = output_filename
            count = 0

            if file_type == "xlsx":
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(file_data))
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.xlsx"

                for sn in wb.sheetnames:
                    ws = wb[sn]
                    for row in ws.iter_rows():
                        for cell in row:
                            if cell.value is None:
                                continue
                            if isinstance(cell.value, str) and find_text in cell.value:
                                cell.value = cell.value.replace(find_text, replace_with)
                                count += 1
                            elif not isinstance(cell.value, str):
                                sval = str(cell.value)
                                if find_text in sval:
                                    cell.value = replace_with
                                    count += 1

                wb.save(out)
                wb.close()

            elif file_type == "xls":
                # Convert .xls to .xlsx, then apply same replace logic
                file_data = _xls_to_xlsx(file_data)
                file_type = "xlsx"
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.xlsx"
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(file_data))

                for sn in wb.sheetnames:
                    ws = wb[sn]
                    for row in ws.iter_rows():
                        for cell in row:
                            if cell.value is None:
                                continue
                            if isinstance(cell.value, str) and find_text in cell.value:
                                cell.value = cell.value.replace(find_text, replace_with)
                                count += 1
                            elif not isinstance(cell.value, str):
                                sval = str(cell.value)
                                if find_text in sval:
                                    cell.value = replace_with
                                    count += 1

                wb.save(out)
                wb.close()

            elif file_type == "docx":
                from docx import Document
                doc = Document(io.BytesIO(file_data))
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.docx"

                for para in doc.paragraphs:
                    if find_text in para.text:
                        # Preserve formatting: replace in runs
                        full_text = para.text
                        if find_text in full_text:
                            new_text = full_text.replace(find_text, replace_with)
                            # Clear all runs and set new text in first run
                            if para.runs:
                                para.runs[0].text = new_text
                                for run in para.runs[1:]:
                                    run.text = ""
                                count += 1
                            else:
                                para.text = new_text
                                count += 1

                # Also replace in tables
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            if find_text in cell.text:
                                for para in cell.paragraphs:
                                    if find_text in para.text:
                                        if para.runs:
                                            para.runs[0].text = para.text.replace(find_text, replace_with)
                                            for run in para.runs[1:]:
                                                run.text = ""
                                        else:
                                            para.text = para.text.replace(find_text, replace_with)
                                        count += 1

                doc.save(out)

            elif file_type == "pptx":
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_data))
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.pptx"

                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame"):
                            for para in shape.text_frame.paragraphs:
                                if find_text in para.text:
                                    if para.runs:
                                        para.runs[0].text = para.text.replace(find_text, replace_with)
                                        for run in para.runs[1:]:
                                            run.text = ""
                                    else:
                                        para.text = para.text.replace(find_text, replace_with)
                                    count += 1

                prs.save(out)

            else:
                return json.dumps({"error": f"Unsupported type: {file_type}"})

            out.seek(0)
            url, name = await self._save_and_link(out.read(), out_name)
            if url:
                return f"[{name}]({url})\n\nReplaced '{find_text}' with '{replace_with}' in {count} place(s), preserving all formatting."
            return json.dumps({"error": "Could not save file"})

        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    # -----------------------------------------------------------------
    # CREATE NEW FILE
    # -----------------------------------------------------------------
    async def create_file(
        self,
        file_type: str,
        content: str,
        output_filename: str = "",
        __user__=None,
        __request__=None,
    ) -> str:
        """Create a new Office file from scratch.

        For xlsx: content is CSV with headers on first line.
        For docx: content is plain text (one paragraph per line).
        For pptx: each line defines a slide. Use "---" as separator between slides.

        Args:
            file_type: 'xlsx', 'docx', or 'pptx'
            content: Content specification
            output_filename: Output filename
        """
        try:
            ftype = file_type.lower().replace(".", "")
            if ftype not in ("xlsx", "docx", "pptx"):
                return json.dumps({"error": f"Unsupported type: {file_type}. Use xlsx, docx, or pptx."})

            out_name = output_filename or f"document.{ftype}"
            out = io.BytesIO()

            if ftype == "xlsx":
                import openpyxl
                from openpyxl.styles import Font, PatternFill, Alignment
                wb = openpyxl.Workbook()
                ws = wb.active

                import csv as _csv_mod
                reader = _csv_mod.reader(io.StringIO(content))
                rows = list(reader)

                if rows:
                    # First row = headers (styled)
                    for j, h in enumerate(rows[0], 1):
                        c = ws.cell(row=1, column=j, value=h.strip())
                        c.font = Font(bold=True, color="FFFFFF", size=11)
                        c.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
                        c.alignment = Alignment(horizontal="center")

                    # Data rows
                    for i, rd in enumerate(rows[1:], 2):
                        for j, v in enumerate(rd, 1):
                            v = v.strip()
                            c = ws.cell(row=i, column=j)
                            if v == '':
                                c.value = None
                            else:
                                try:
                                    c.value = int(v)
                                except ValueError:
                                    try:
                                        c.value = float(v)
                                    except ValueError:
                                        if v.lower() == 'true':
                                            c.value = True
                                        elif v.lower() == 'false':
                                            c.value = False
                                        else:
                                            c.value = v
                            c.alignment = Alignment(
                                horizontal="center" if isinstance(c.value, (int, float)) else "left"
                            )

                    # Auto-fit columns
                    for col in ws.columns:
                        mx = max((len(str(c.value or "")) for c in col), default=5)
                        ws.column_dimensions[col[0].column_letter].width = min(mx + 3, 50)

                wb.save(out)
                wb.close()

            elif ftype == "docx":
                from docx import Document
                from docx.shared import Pt
                doc = Document()
                for line in content.split("\n"):
                    doc.add_paragraph(line)
                doc.save(out)

            elif ftype == "pptx":
                from pptx import Presentation
                from pptx.util import Inches
                prs = Presentation()
                slide_specs = re.split(r'\n---\n|\r\n---\r\n|\n---\n', content)
                blank_layout = prs.slide_layouts[6]

                for spec in slide_specs:
                    spec = spec.strip()
                    if not spec:
                        continue
                    lines = spec.split("\n")
                    title = lines[0].strip()
                    body = "\n".join(lines[1:]).strip() if len(lines) > 1 else ""

                    slide = prs.slides.add_slide(blank_layout)
                    txBox = slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.3), Inches(9), Inches(1)
                    )
                    tf = txBox.text_frame
                    tf.text = title
                    p = tf.paragraphs[0]
                    p.font.size = Inches(0.6)
                    p.font.bold = True

                    if body:
                        txBox2 = slide.shapes.add_textbox(
                            Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)
                        )
                        tf2 = txBox2.text_frame
                        tf2.text = body
                        for para in tf2.paragraphs:
                            para.font.size = Inches(0.3)

                prs.save(out)

            out.seek(0)
            url, name = await self._save_and_link(out.read(), out_name)
            if url:
                return f"[{name}]({url})\n\nCreated new {ftype.upper()} file."
            return json.dumps({"error": "Could not save file"})

        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})


    async def generate_document(self, content: str, template: str = "", accent: str = "", output_filename: str = "", __user__=None, __request__=None) -> str:
        """Generate a professional .docx document with optional cover page, callouts, signatures, headers/footers, and styled headings.
        Uses YAML front matter (--- delimited) for metadata: title, author, header, footer, cover (page).
        Body supports markdown headings (# ## ###), callout blocks (::: callout type="info" title="..."), signatures (::: signature name="..." role="..."), and [[toc]] placeholder.
        Template 'report' enables cover page automatically. Accent color applied to headings and title."""
        try:
            from docx import Document
            from docx.shared import Pt, Inches, Cm, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            import re as _re
            
            acc = accent or getattr(self.valves, 'default_accent', '#1B6B93')
            r, g, b = int(acc[1:3], 16), int(acc[3:5], 16), int(acc[5:7], 16)
            
            doc = Document()
            style = doc.styles['Normal']
            style.font.name = 'Calibri'
            style.font.size = Pt(11)
            
            parsed = content
            title = "Document"
            author_name = getattr(self.valves, 'default_author', '')
            header_text = ""
            footer_text = ""
            has_cover = False
            
            if content.startswith('---'):
                parts = content.split('---', 2)
                if len(parts) >= 3:
                    yaml_text = parts[1]
                    parsed = parts[2]
                    for line in yaml_text.strip().split('\n'):
                        if ':' in line:
                            k, v = line.split(':', 1)
                            k, v = k.strip().lower(), v.strip()
                            if k == 'title': title = v
                            elif k == 'cover' and v.lower() != 'none': has_cover = True
                            elif k == 'author': author_name = v
                            elif k == 'header': header_text = v
                            elif k == 'footer': footer_text = v
            
            if has_cover or template == 'report':
                for _ in range(8):
                    doc.add_paragraph()
                p = doc.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run = p.add_run(title)
                run.font.size = Pt(28)
                run.font.color.rgb = RGBColor(r, g, b)
                run.bold = True
                if author_name:
                    p3 = doc.add_paragraph()
                    p3.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    p3.add_run(f"By {author_name}").font.size = Pt(11)
                doc.add_page_break()
            
            if header_text:
                hp = doc.sections[0].header.paragraphs[0]
                hp.text = header_text
                hp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            if footer_text:
                fp = doc.sections[0].footer.paragraphs[0]
                fp.text = footer_text
                fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            for raw in parsed.strip().split('\n'):
                line = raw.strip()
                if line.startswith('```'):
                    continue
                if line == '[[toc]]':
                    p = doc.add_paragraph('[Table of Contents - update in Word]')
                    p.runs[0].italic = True
                    continue
                h_match = _re.match(r'^(#{1,3})\s+(.*)', line)
                if h_match:
                    lvl = len(h_match.group(1))
                    heading = doc.add_heading(h_match.group(2), level=lvl)
                    heading.runs[0].font.color.rgb = RGBColor(r, g, b) if lvl <= 2 else RGBColor(0, 0, 0)
                    continue
                if line.startswith('::: callout'):
                    m = _re.match(r'::: callout type="(\w+)" title="(.+?)"', line)
                    if m:
                        colors = {'info': (33,150,243), 'success': (76,175,80), 'warning': (255,152,0), 'danger': (244,67,54)}
                        cc = colors.get(m.group(1), (33,150,243))
                        p = doc.add_paragraph()
                        p.paragraph_format.left_indent = Cm(0.5)
                        r1 = p.add_run(f"{m.group(1).upper()}: {m.group(2)}")
                        r1.bold = True
                        r1.font.color.rgb = RGBColor(*cc)
                        continue
                if line.startswith('::: signature'):
                    m = _re.match(r'::: signature name="(.+?)"\s*(?:role="(.+?)")?', line)
                    if m:
                        p = doc.add_paragraph()
                        p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                        p.add_run(m.group(1)).bold = True
                        if m.group(2):
                            p2 = doc.add_paragraph()
                            p2.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                            p2.add_run(m.group(2))
                        continue
                if line.startswith(':::') or line == ':::':
                    continue
                if line:
                    doc.add_paragraph(line)
            
            out = io.BytesIO()
            doc.save(out)
            out.seek(0)
            fname = output_filename or f"{title.replace(' ', '_')}.docx"
            url, name = self._save_and_link(out.read(), fname)
            if url:
                return f"[{name}]({url})\n\nProfessional document generated."
            return json.dumps({"error": "Could not save file"})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    async def generate_slides(self, content: str, theme: str = "", accent: str = "", output_filename: str = "", __user__=None, __request__=None) -> str:
        """Generate a professional .pptx presentation from JSON slide spec or plain text.
        JSON format: {"slides": [{"layout": "cover|title_bullets|kpi_row|chart|closing", "title": "...", ...}]}
        Layouts: cover (title+subtitle), title_bullets (title + body/bullets), kpi_row (stat cards with value/label/change), chart (bar chart from labels/values), closing (title + takeaways).
        Themes: midnight, charcoal, slate, coral, forest, ocean (or custom accent).
        If content is plain text, auto-wraps as a title_bullets slide."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            from pptx.chart.data import CategoryChartData
            
            acc = accent or getattr(self.valves, 'default_accent', '#1B6B93')
            r, g, b = int(acc[1:3], 16), int(acc[3:5], 16), int(acc[5:7], 16)
            thm = theme or getattr(self.valves, 'default_theme', 'auto')
            
            slides_data = []
            try:
                spec = json.loads(content)
                slides_data = spec.get('slides', [])
            except Exception:
                slides_data = [{'layout': 'title_bullets', 'title': 'Presentation', 'body': content}]
            
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            theme_colors = {
                'midnight': (18, 18, 36),
                'charcoal': (30, 30, 30),
                'slate': (40, 50, 60),
                'coral': (255, 100, 80),
                'forest': (30, 80, 40),
                'ocean': (20, 60, 100)
            }
            bg = theme_colors.get(thm, (18, 18, 36))
            
            for ss in slides_data:
                layout = ss.get('layout', 'title_bullets')
                stitle = ss.get('title', 'Slide')
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(*bg)
                
                if layout == 'cover':
                    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
                    p = tb.text_frame.paragraphs[0]
                    p.text = stitle
                    p.font.size = Pt(44)
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER
                    iu = ss.get('image_url',ss.get('image_hint',''))
                    if iu:
                        tb3 = slide.shapes.add_textbox(Inches(2),Inches(6),Inches(9),Inches(0.5))
                        p3 = tb3.text_frame.paragraphs[0]
                        p3.text = 'Image: '+iu
                        p3.font.size = Pt(10)
                        p3.font.color.rgb = RGBColor(150,150,150)
                        p3.font.italic = True
                    if 'subtitle' in ss:
                        tb2 = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(9), Inches(1))
                        p2 = tb2.text_frame.paragraphs[0]
                        p2.text = ss['subtitle']
                        p2.font.size = Pt(18)
                        p2.font.color.rgb = RGBColor(200, 200, 200)
                        p2.alignment = PP_ALIGN.CENTER
                
                elif layout == 'kpi_row':
                    stats = ss.get('stats', [])
                    for idx, stat in enumerate(stats[:4]):
                        shape = slide.shapes.add_shape(1, Inches(0.5 + idx * 3.2), Inches(2.5), Inches(3), Inches(2.5))
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                        shape.line.fill.background()
                        tf = shape.text_frame
                        p = tf.paragraphs[0]
                        p.text = str(stat.get('value', ''))
                        p.font.size = Pt(32)
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(r, g, b)
                        p.alignment = PP_ALIGN.CENTER
                        p2 = tf.add_paragraph()
                        p2.text = stat.get('label', '')
                        p2.font.size = Pt(12)
                        p2.font.color.rgb = RGBColor(100, 100, 100)
                        p2.alignment = PP_ALIGN.CENTER
                        if 'change' in stat:
                            p3 = tf.add_paragraph()
                            p3.text = stat['change']
                            p3.font.size = Pt(14)
                            p3.font.color.rgb = RGBColor(76, 175, 80) if '+' in stat['change'] else RGBColor(244, 67, 54)
                            p3.alignment = PP_ALIGN.CENTER
                
                elif layout == 'chart':
                    labels = ss.get('labels', [])
                    values = ss.get('values', [])
                    cd = CategoryChartData()
                    cd.categories = labels
                    cd.add_series('Values', values)
                    _chart_types = {
                        'column': 2, 'bar': 51, 'line': 4, 'pie': 5,
                        'area': 76, 'scatter': 74, 'radar': 10,
                        'doughnut': -4121, 'bubble': 15, 'stock': 88,
                    }
                    _ct = ss.get('chart_type', 'column')
                    _chart_enum = _chart_types.get(_ct, 2)
                    chart = slide.shapes.add_chart(_chart_enum, Inches(1), Inches(1.5), Inches(11), Inches(5), cd).chart
                    chart.has_legend = False
                
                elif layout == 'title_bullets':
                    tb = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1.2))
                    p = tb.text_frame.paragraphs[0]
                    p.text = stitle
                    p.font.size = Pt(36)
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.font.bold = True
                    body = ss.get('body', ss.get('bullets', ''))
                    if isinstance(body, list):
                        body = '\n'.join(body)
                    if body:
                        tb2 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(5))
                        p2 = tb2.text_frame.paragraphs[0]
                        p2.text = body
                        p2.font.size = Pt(18)
                        p2.font.color.rgb = RGBColor(220, 220, 220)
                
                elif layout == 'closing':
                    tb = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9), Inches(2))
                    p = tb.text_frame.paragraphs[0]
                    p.text = stitle
                    p.font.size = Pt(40)
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER
                    for t in ss.get('takeaways', []):
                        p2 = tb.text_frame.add_paragraph()
                        p2.text = f"  {t}"
                        p2.font.size = Pt(16)
                        p2.font.color.rgb = RGBColor(180, 180, 180)
                        p2.alignment = PP_ALIGN.CENTER

                elif layout == 'section':
                    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(3))
                    p = tb.text_frame.paragraphs[0]
                    num = ss.get('number', '')
                    p.text = f"{num}  {stitle}"
                    p.font.size = Pt(36)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.alignment = PP_ALIGN.CENTER
                    sub = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11), Inches(1.5))
                    ps = sub.text_frame.paragraphs[0]
                    ps.text = ss.get('subtitle', '')
                    ps.font.size = Pt(18)
                    ps.font.color.rgb = RGBColor(200, 200, 200)
                    ps.alignment = PP_ALIGN.CENTER

                elif layout == 'timeline':
                    items = ss.get('items', [])
                    for idx, item in enumerate(items):
                        y = 1.5 + idx * 1.2
                        tb = slide.shapes.add_textbox(Inches(1.5), Inches(y), Inches(10), Inches(1))
                        p = tb.text_frame.paragraphs[0]
                        time_str = item.get('time', '')
                        label = item.get('label', item.get('title', ''))
                        p.text = f"{time_str}  |  {label}"
                        p.font.size = Pt(18)
                        p.font.color.rgb = RGBColor(255, 255, 255)
                        if item.get('highlight', False):
                            p.font.bold = True

                elif layout == 'icon_grid_2x2':
                    items = ss.get('items', [])[:4]
                    positions = [(1, 1.5), (6.5, 1.5), (1, 4.5), (6.5, 4.5)]
                    for idx, item in enumerate(items):
                        if idx >= len(positions):
                            break
                        x, y = positions[idx]
                        icon = item.get('icon', '')
                        title = item.get('title', '')
                        desc = item.get('description', '')
                        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(5), Inches(2.5))
                        tf = tb.text_frame
                        tf.word_wrap = True
                        p = tf.paragraphs[0]
                        p.text = f"{icon}  {title}"
                        p.font.size = Pt(22)
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(255, 255, 255)
                        if desc:
                            p2 = tf.add_paragraph()
                            p2.text = desc
                            p2.font.size = Pt(14)
                            p2.font.color.rgb = RGBColor(180, 180, 180)

                elif layout == 'pillars':
                    items = ss.get('items', [])
                    n = len(items)
                    if n > 0:
                        col_w = 10.0 / n
                        for idx, item in enumerate(items):
                            x = 1.5 + idx * (col_w + 0.5) if n <= 4 else 1 + idx * (11 / n)
                            tb = slide.shapes.add_textbox(Inches(x), Inches(2), Inches(col_w - 0.3), Inches(4))
                            tf = tb.text_frame
                            tf.word_wrap = True
                            p = tf.paragraphs[0]
                            p.text = item.get('icon', '')
                            p.font.size = Pt(36)
                            p.alignment = PP_ALIGN.CENTER
                            p2 = tf.add_paragraph()
                            p2.text = item.get('title', '')
                            p2.font.size = Pt(16)
                            p2.font.bold = True
                            p2.font.color.rgb = RGBColor(255, 255, 255)
                            p2.alignment = PP_ALIGN.CENTER
                            desc = item.get('description', '')
                            if desc:
                                p3 = tf.add_paragraph()
                                p3.text = desc
                                p3.font.size = Pt(12)
                                p3.font.color.rgb = RGBColor(180, 180, 180)
                                p3.alignment = PP_ALIGN.CENTER

                elif layout == 'quote':
                    q = ss.get('quote', '')
                    author = ss.get('author', '')
                    tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(11), Inches(2.5))
                    tf = tb.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = f"\u201c{q}\u201d"
                    p.font.size = Pt(28)
                    p.font.italic = True
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.alignment = PP_ALIGN.CENTER
                    if author:
                        p2 = tf.add_paragraph()
                        p2.text = f"-- {author}"
                        p2.font.size = Pt(16)
                        p2.font.color.rgb = RGBColor(180, 180, 180)
                        p2.alignment = PP_ALIGN.CENTER

                elif layout == 'alert':
                    lv = ss.get('level','info')
                    cl = {'info':(33,150,243),'warning':(255,152,0),'danger':(244,67,54),'success':(76,175,80)}
                    ac = cl.get(lv,(33,150,243))
                    s = slide.shapes.add_shape(1,Inches(0),Inches(0),Inches(0.2),Inches(7.5))
                    s.fill.solid()
                    s.fill.fore_color.rgb = RGBColor(*ac)
                    s.line.fill.background()
                    tb = slide.shapes.add_textbox(Inches(1),Inches(1),Inches(11),Inches(1))
                    p = tb.text_frame.paragraphs[0]
                    p.text = lv.upper()+': '+stitle
                    p.font.size = Pt(32)
                    p.font.color.rgb = RGBColor(*ac)
                    p.font.bold = True
                    bd = ss.get('body','')
                    if bd:
                        tb2 = slide.shapes.add_textbox(Inches(1),Inches(2.5),Inches(11),Inches(4))
                        p2 = tb2.text_frame.paragraphs[0]
                        p2.text = bd
                        p2.font.size = Pt(18)
                        p2.font.color.rgb = RGBColor(200,200,200)

                elif layout == 'table':
                    rows = ss.get('rows', [])
                    headers = ss.get('headers', [])
                    if headers and rows:
                        n_cols = len(headers)
                        n_rows = len(rows) + 1
                        tbl_x, tbl_y = Inches(1), Inches(1.5)
                        tbl_w, tbl_h = Inches(12), Inches(0.5 * n_rows)
                        table_shape = slide.shapes.add_table(n_rows, n_cols, tbl_x, tbl_y, tbl_w, tbl_h)
                        table = table_shape.table
                        for ci, h in enumerate(headers):
                            cell = table.cell(0, ci)
                            cell.text = str(h)
                            for paragraph in cell.text_frame.paragraphs:
                                paragraph.font.size = Pt(14)
                                paragraph.font.bold = True
                                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                        for ri, row in enumerate(rows):
                            for ci in range(n_cols):
                                cell = table.cell(ri + 1, ci)
                                cell.text = str(row[ci]) if ci < len(row) else ''
                                for paragraph in cell.text_frame.paragraphs:
                                    paragraph.font.size = Pt(12)
                                    paragraph.font.color.rgb = RGBColor(220, 220, 220)

                elif layout == 'title_bar':
                    tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1))
                    p = tb.text_frame.paragraphs[0]
                    p.text = stitle
                    p.font.size = Pt(40)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.alignment = PP_ALIGN.CENTER
                    sub = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1))
                    ps = sub.text_frame.paragraphs[0]
                    ps.text = ss.get('subtitle', t)
                    ps.font.size = Pt(18)
                    ps.font.color.rgb = RGBColor(200, 200, 200)
                    ps.alignment = PP_ALIGN.CENTER            
            out = io.BytesIO()
            prs.save(out)
            out.seek(0)
            fname = output_filename or "presentation.pptx"
            url, name = self._save_and_link(out.read(), fname)
            if url:
                return f"[{name}]({url})\n\nPresentation with {len(slides_data)} slides."
            return json.dumps({"error": "Could not save file"})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})


    async def generate_spreadsheet(self, content: str, template: str = "", accent: str = "", output_filename: str = "", __user__=None, __request__=None) -> str:
        """Generate a professional Excel workbook from a JSON specification.

        Accepts a JSON spec with title, sheets array. Each sheet has a kind:
        table, inputs, kpi_row, chart, notes. Supports Excel tables with
        auto-filter, row stripes, freeze panes, column formatting (text,
        number, currency, percent, date, integer), live formulas, conditional
        formatting, data validation, and multi-sheet workbooks.

        Args:
            content: JSON specification string defining the workbook structure
            template: Template style ('financial', 'minimal', or empty for default)
            accent: Hex color accent (e.g. '#1B6B93'), defaults to self.valves.default_accent
            output_filename: Custom output filename (defaults to title-based .xlsx)
        Returns:
            Markdown link to the generated file, or error JSON
        """
        try:
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side, numbers
            from openpyxl.utils import get_column_letter
            from openpyxl.worksheet.table import Table, TableStyleInfo
            from openpyxl.formatting.rule import DataBarRule, ColorScaleRule, CellIsRule
            from openpyxl.worksheet.datavalidation import DataValidation

            acc = accent or getattr(self.valves, 'default_accent', '#1B6B93')
            r, g, b = int(acc[1:3], 16), int(acc[3:5], 16), int(acc[5:7], 16)
            acc_hex = acc.lstrip('#')

            spec = json.loads(content)
            title = spec.get('title', 'Workbook')
            sheets_spec = spec.get('sheets', [])

            wb = openpyxl.Workbook()
            wb.remove(wb.active)

            # Color palette
            header_fill = PatternFill(start_color=acc_hex, end_color=acc_hex, fill_type="solid")
            header_font = Font(bold=True, color="FFFFFF", size=11)
            alt_fill = PatternFill(start_color="F5F5F5", end_color="F5F5F5", fill_type="solid")
            thin_border = Border(
                left=Side(style='thin', color='D0D0D0'),
                right=Side(style='thin', color='D0D0D0'),
                top=Side(style='thin', color='D0D0D0'),
                bottom=Side(style='thin', color='D0D0D0')
            )
            number_fmts = {
                'currency': '#,##0.00',
                'percent': '0.0%',
                'number': '#,##0',
                'date': 'YYYY-MM-DD',
                'text': '@',
                'integer': '#,##0'
            }

            # Template-specific styles
            input_font = None
            input_fill = None
            if template == 'financial':
                input_font = Font(color="1F4E79", size=11)
                input_fill = PatternFill(start_color="FFF2CC", end_color="FFF2CC", fill_type="solid")

            for sheet_spec in sheets_spec:
                name = sheet_spec.get('name', 'Sheet')
                kind = sheet_spec.get('kind', 'table')

                ws = wb.create_sheet(title=name)

                if kind == 'table':
                    table_spec = sheet_spec.get('table', {})
                    columns = table_spec.get('columns', [])
                    rows = table_spec.get('rows', [])

                    # Write headers
                    for j, col in enumerate(columns, 1):
                        cell = ws.cell(row=1, column=j, value=col.get('header', col.get('key', '')))
                        cell.font = header_font
                        cell.fill = header_fill
                        cell.alignment = Alignment(horizontal='center')
                        cell.border = thin_border

                    # Write data rows
                    for i, row in enumerate(rows, 2):
                        for j, col in enumerate(columns, 1):
                            key = col['key']
                            val = row.get(key, '')
                            cell = ws.cell(row=i, column=j)

                            # Handle formulas (strings starting with =)
                            if isinstance(val, str) and val.startswith('='):
                                cell.value = val
                                if template == 'financial':
                                    cell.font = Font(color="000000", size=11)
                            else:
                                cell.value = val
                                if template == 'financial' and col.get('format') in ('currency', 'number'):
                                    cell.font = Font(color="1F4E79", size=11)
                                    if input_fill:
                                        cell.fill = input_fill

                            # Number format
                            fmt = col.get('format', '')
                            if fmt in number_fmts:
                                cell.number_format = number_fmts[fmt]

                            cell.border = thin_border
                            cell.alignment = Alignment(horizontal='right' if fmt in ('currency', 'number', 'percent', 'integer') else 'left')

                            # Zebra striping
                            if i % 2 == 0:
                                cell.fill = alt_fill

                    # Freeze panes
                    freeze = table_spec.get('freeze', '')
                    if freeze:
                        ws.freeze_panes = freeze

                    # Auto-fit columns
                    for j, col in enumerate(columns, 1):
                        max_len = 0
                        col_letter = get_column_letter(j)
                        for i in range(1, min(len(rows) + 2, 100)):
                            val = str(ws.cell(row=i, column=j).value or '')
                            max_len = max(max_len, len(val))
                        ws.column_dimensions[col_letter].width = min(max_len + 3, 40)

                    # Excel Table (ListObject) with auto-filter and styling
                    if table_spec.get('excel_table', False) and rows:
                        max_row = len(rows) + 1
                        max_col = len(columns)
                        ref = f"A1:{get_column_letter(max_col)}{max_row}"
                        tab = Table(displayName=table_spec.get('name', 'Table1').replace(' ', '_'), ref=ref)
                        style = TableStyleInfo(
                            name="TableStyleMedium6",
                            showFirstColumn=False,
                            showLastColumn=False,
                            showRowStripes=True,
                            showColumnStripes=False
                        )
                        tab.tableStyleInfo = style
                        ws.add_table(tab)

                    # Auto-filter (when not using excel_table)
                    if not table_spec.get('excel_table', False) and rows and table_spec.get('auto_filter', False):
                        max_row = len(rows) + 1
                        max_col = len(columns)
                        ws.auto_filter.ref = f"A1:{get_column_letter(max_col)}{max_row}"

                    # Conditional formatting
                    cf = table_spec.get('conditional_formatting', {})
                    if cf.get('type') == 'data_bar':
                        col_range = cf.get('range', 'B2:B100')
                        rule = DataBarRule(start_type='min', end_type='max', color=acc_hex)
                        ws.conditional_formatting.add(col_range, rule)
                    elif cf.get('type') == 'color_scale':
                        col_range = cf.get('range', 'B2:B100')
                        rule = ColorScaleRule(
                            start_type='min', start_color='FFFFFF',
                            end_type='max', end_color=acc_hex
                        )
                        ws.conditional_formatting.add(col_range, rule)
                    elif cf.get('type') == 'cell_is':
                        col_range = cf.get('range', 'B2:B100')
                        rule = CellIsRule(
                            operator=cf.get('operator', 'greaterThan'),
                            formula=[cf.get('value', '0')],
                            fill=PatternFill(start_color=cf.get('fill_color', 'C6EFCE'),
                                             end_color=cf.get('fill_color', 'C6EFCE'),
                                             fill_type='solid')
                        )
                        ws.conditional_formatting.add(col_range, rule)

                    # Data Validation
                    validation = table_spec.get('validation', {})
                    if validation:
                        dv = DataValidation(
                            type=validation.get('type', 'list'),
                            formula1=validation.get('formula', ''),
                            allow_blank=True
                        )
                        dv.error = "Invalid value"
                        dv.errorTitle = "Validation Error"
                        ws.add_data_validation(dv)
                        col_range = validation.get('range', 'A2:A1000')
                        dv.add(col_range)

                elif kind == 'kpi_row':
                    stats = sheet_spec.get('stats', [])
                    title_val = sheet_spec.get('title', '')
                    if title_val:
                        ws.cell(row=1, column=1, value=title_val).font = Font(bold=True, size=14, color=acc_hex)
                    for idx, stat in enumerate(stats):
                        col = idx * 3 + 1
                        ws.cell(row=3, column=col, value=stat.get('value', '')).font = Font(bold=True, size=24, color="000000")
                        ws.cell(row=4, column=col, value=stat.get('label', '')).font = Font(size=11, color="808080")
                        change = stat.get('change', '')
                        if change:
                            c = ws.cell(row=5, column=col, value=change)
                            c.font = Font(size=12, color="27AE60" if '+' in str(change) else "E74C3C")
                    # Auto-fit KPI columns
                    for idx in range(len(stats)):
                        col_letter = get_column_letter(idx * 3 + 1)
                        ws.column_dimensions[col_letter].width = 18

                elif kind == 'notes':
                    note_spec = sheet_spec
                    colors = {'info': '1F4E79', 'success': '27AE60', 'warning': 'E67E22', 'danger': 'E74C3C'}
                    note_color = colors.get(note_spec.get('level', 'info'), '1F4E79')
                    title_val = note_spec.get('title', '')
                    if title_val:
                        ws.cell(row=1, column=1, value=title_val).font = Font(bold=True, size=14, color=note_color)
                    note_text = note_spec.get('text', '')
                    if note_text:
                        ws.cell(row=3, column=1, value=note_text).font = Font(size=11)
                    ws.column_dimensions['A'].width = 80

                elif kind == 'inputs':
                    items = sheet_spec.get('items', [])
                    title_val = sheet_spec.get('title', '')
                    if title_val:
                        ws.cell(row=1, column=1, value=title_val).font = Font(bold=True, size=14, color=acc_hex)
                    for idx, item in enumerate(items):
                        row = idx * 3 + 3
                        label_cell = ws.cell(row=row, column=1, value=item.get('label', ''))
                        label_cell.font = Font(bold=True, size=11)
                        val_cell = ws.cell(row=row, column=2, value=item.get('value', ''))
                        if template == 'financial' and input_fill:
                            val_cell.fill = input_fill
                            val_cell.font = input_font or Font(color="1F4E79", size=11)
                        unit = item.get('unit', '')
                        if unit:
                            ws.cell(row=row, column=3, value=unit).font = Font(size=10, color="808080")
                        comment = item.get('comment', '')
                        if comment:
                            ws.cell(row=row+1, column=1, value=comment).font = Font(size=9, color="808080", italic=True)
                    ws.column_dimensions['A'].width = 20
                    ws.column_dimensions['B'].width = 15

                elif kind == 'chart':
                    # Placeholder: openpyxl chart support
                    chart_type = sheet_spec.get('chart_type', 'bar')
                    chart_title = sheet_spec.get('chart_title', name)
                    data_ref = sheet_spec.get('data_ref', 'A1:B10')
                    ws.cell(row=1, column=1, value=f"Chart: {chart_title}").font = Font(bold=True, size=14, color=acc_hex)
                    ws.cell(row=2, column=1, value=f"Type: {chart_type}").font = Font(size=11, color="808080")
                    ws.cell(row=3, column=1, value=f"Data reference: {data_ref}").font = Font(size=11, color="808080")
                    ws.cell(row=5, column=1, value="Note: Chart rendering requires data from a table sheet.").font = Font(size=10, color="808080", italic=True)

            out = io.BytesIO()
            wb.save(out)
            out.seek(0)
            fname = output_filename or f"{title.replace(' ', '_')}.xlsx"
            url, name = self._save_and_link(out.read(), fname)
            if url:
                return f"[{name}]({url})\n\nProfessional workbook with {len(sheets_spec)} sheets generated."
            return json.dumps({"error": "Could not save file"})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})


    async def tracked_change(self, file_id: str, change_type: str, content: str, author: str = "Reviewer", paragraph_index: int = -1, output_filename: str = "", __user__=None, __request__=None) -> str:
        """Apply tracked changes (redlines) to a Word document with custom author name.
    
        change_type: replace (use old_text|||new_text), insert (append text with redline), delete (mark paragraph for deletion)
        author: Name shown in Word's Track Changes (e.g., "Sergio Pedro")
        """
        try:
            import sqlite3 as s3
            conn2 = s3.connect(r"C:\\Users\\Administrator\\AppData\\Roaming\\open-webui\\data\\webui.db")
            row = conn2.execute("SELECT filename, meta FROM file WHERE id=?", (file_id,)).fetchone()
            if not row:
                row = conn2.execute("SELECT filename, meta FROM file WHERE filename LIKE ?", (f"%{file_id}%",)).fetchone()
            if not row:
                conn2.close()
                return json.dumps({"error": "File not found"})
            filename = row[0]
            meta = json.loads(row[1]) if row[1] else {}
            fp = meta.get("path", file_id)
            if not os.path.exists(fp):
                fp = os.path.join(os.environ.get("APPDATA",""), "open-webui", "data", "uploads", os.path.basename(fp))
            if not os.path.exists(fp):
                conn2.close()
                return json.dumps({"error": "File not found on disk"})
            with open(fp, "rb") as f:
                data = f.read()
            conn2.close()
    
            from docx import Document
            from docx_revisions import RevisionParagraph
            doc = Document(io.BytesIO(data))
            out_name = output_filename or filename
            results = []
    
            if change_type == "replace":
                parts = content.split("|||", 1)
                if len(parts) != 2:
                    return json.dumps({"error": "Format: old_text|||new_text"})
                find_t, replace_t = parts
                for i, p in enumerate(doc.paragraphs):
                    if paragraph_index >= 0 and i != paragraph_index:
                        continue
                    if find_t in p.text:
                        rp = RevisionParagraph.from_paragraph(p)
                        cnt = rp.replace_tracked(find_t, replace_t, author=author)
                        results.append(f"Para {i}: {cnt} replacements")
            elif change_type == "insert":
                p = doc.add_paragraph()
                rp = RevisionParagraph.from_paragraph(p)
                rp.add_tracked_insertion(content, author=author)
                results.append("Inserted tracked text")
            elif change_type == "delete":
                idx = int(content) if content.isdigit() else paragraph_index
                if idx >= 0 and idx < len(doc.paragraphs):
                    p = doc.paragraphs[idx]
                    rp = RevisionParagraph.from_paragraph(p)
                    rp.add_tracked_deletion(0, len(p.text), author=author)
                    results.append(f"Marked para {idx} for deletion")
    
            out = io.BytesIO()
            doc.save(out)
            out.seek(0)
            url, name = self._save_and_link(out.read(), out_name)
            if url:
                return f"[{name}]({url})\n\nTracked changes by '{author}':\n" + "\n".join(results)
            return json.dumps({"error": "Could not save file"})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})


    async def merge_sheets(self, file_ids: str, output_filename: str = "", __user__=None, __request__=None) -> str:
        try:
            import sqlite3 as s3, openpyxl, io, os
            from copy import copy
            conn2 = s3.connect(r"C:\\Users\\Administrator\\AppData\\Roaming\\open-webui\\data\\webui.db")
            ids = [fid.strip() for fid in file_ids.split(",") if fid.strip()]
            wb_out = openpyxl.Workbook()
            wb_out.remove(wb_out.active)
            merged = 0
            for fid in ids:
                row = conn2.execute("SELECT filename, meta FROM file WHERE id=?", (fid,)).fetchone()
                if not row:
                    row = conn2.execute("SELECT filename, meta FROM file WHERE filename LIKE ?", ("%"+fid+"%",)).fetchone()
                if not row:
                    continue
                filename = row[0]
                meta = json.loads(row[1]) if row[1] else {}
                fp = meta.get("path", fid)
                if not os.path.exists(fp):
                    alt = os.path.join(os.environ.get("APPDATA",""), "open-webui","data","uploads", os.path.basename(fp))
                    fp = alt if os.path.exists(alt) else ""
                if not fp or not os.path.exists(fp):
                    continue
                wb_src = openpyxl.load_workbook(io.BytesIO(open(fp,"rb").read()))
                base_name = os.path.splitext(os.path.basename(filename))[0][:15]
                for sn in wb_src.sheetnames:
                    ws_src = wb_src[sn]
                    sheet_name = (base_name + "_" + sn)[:31]
                    ws_out = wb_out.create_sheet(title=sheet_name)
                    for ri, row_data in enumerate(ws_src.iter_rows(), 1):
                        for ci, cell in enumerate(row_data, 1):
                            out_cell = ws_out.cell(row=ri, column=ci, value=cell.value)
                            if cell.has_style:
                                out_cell.font = copy(cell.font)
                                out_cell.fill = copy(cell.fill)
                                out_cell.border = copy(cell.border)
                                out_cell.alignment = copy(cell.alignment)
                                out_cell.number_format = cell.number_format
                    merged += 1
                wb_src.close()
            conn2.close()
            if merged == 0:
                return json.dumps({"error": "No files could be merged"})
            out = io.BytesIO()
            wb_out.save(out)
            out.seek(0)
            fname = output_filename or "merged_workbook.xlsx"
            url, name = self._save_and_link(out.read(), fname)
            if url:
                return f"[{name}]({url})\n\nMerged {merged} sheets from {len(ids)} files."
            return json.dumps({"error": "Could not save file"})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    async def batch_process(self, file_ids: str, operation: str, params: str = "", output_filename: str = "", __user__=None, __request__=None) -> str:
        try:
            ids = [fid.strip() for fid in file_ids.split(",") if fid.strip()]
            results = []
            for fid in ids:
                if operation == "replace":
                    parts = params.split("|||", 1)
                    if len(parts) == 2:
                        await self.replace_text(fid, parts[0], parts[1], "", __user__, __request__)
                        results.append(f"  {fid}: replaced")
                elif operation == "add_rows":
                    await self.add_content(fid, params, "", __user__, __request__)
                    results.append(f"  {fid}: rows added")
            if results:
                return "Batch processed " + str(len(ids)) + " files:\n" + "\n".join(results)
            return json.dumps({"error": "No files processed"})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    async def auto_backup(self, __user__=None, __request__=None) -> str:
        try:
            import shutil, datetime
            db_path = r"C:\\Users\\Administrator\\AppData\\Roaming\\open-webui\\data\\webui.db"
            backup_dir = os.path.join(os.environ.get("LOCALAPPDATA", os.path.expanduser("~")), "open-webui", "backups")
            os.makedirs(backup_dir, exist_ok=True)
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            backup_name = f"webui_backup_{timestamp}.db"
            backup_path = os.path.join(backup_dir, backup_name)
            shutil.copy2(db_path, backup_path)
            size_kb = os.path.getsize(backup_path) / 1024
            return json.dumps({"success": True, "backup_path": backup_path, "size_kb": round(size_kb,1), "message": f"Backup: {backup_name} ({size_kb:.1f} KB)"})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})



    async def merge_pdfs(self, file_ids: str, output_filename: str = "", __user__=None, __request__=None) -> str:
        try:
            import fitz, sqlite3 as s3, io, os
            conn2 = s3.connect(r"C:\\Users\\Administrator\\AppData\\Roaming\\open-webui\\data\\webui.db")
            ids = [fid.strip() for fid in file_ids.split(",") if fid.strip()]
            merger = fitz.open()
            count = 0
            for fid in ids:
                row = conn2.execute("SELECT meta FROM file WHERE id=?", (fid,)).fetchone()
                if not row:
                    row = conn2.execute("SELECT meta FROM file WHERE filename LIKE ?", ("%"+fid+"%",)).fetchone()
                if not row:
                    continue
                meta = json.loads(row[0]) if row[0] else {}
                fp = meta.get("path", fid)
                if not os.path.exists(fp):
                    fp = os.path.join(os.environ.get("APPDATA",""), "open-webui","data","uploads", os.path.basename(fp))
                if not os.path.exists(fp):
                    continue
                src = fitz.open(fp)
                merger.insert_pdf(src)
                src.close()
                count += 1
            conn2.close()
            if count == 0:
                merger.close()
                return json.dumps({"error": "No PDFs could be merged"})
            out = io.BytesIO()
            merger.save(out)
            merger.close()
            out.seek(0)
            fname = output_filename or "merged.pdf"
            url, name = self._save_and_link(out.read(), fname)
            if url:
                return f"[{name}]({url})\n\nMerged {count} PDFs into one file."
            return json.dumps({"error": "Could not save file"})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    async def split_pdf(self, file_id: str, pages_per_file: int = 1, output_filename: str = "", __user__=None, __request__=None) -> str:
        try:
            import fitz, sqlite3 as s3, io, os
            conn2 = s3.connect(r"C:\\Users\\Administrator\\AppData\\Roaming\\open-webui\\data\\webui.db")
            row = conn2.execute("SELECT meta FROM file WHERE id=?", (file_id,)).fetchone()
            if not row:
                row = conn2.execute("SELECT meta FROM file WHERE filename LIKE ?", ("%"+file_id+"%",)).fetchone()
            if not row:
                conn2.close()
                return json.dumps({"error": "File not found"})
            meta = json.loads(row[0]) if row[0] else {}
            fp = meta.get("path", file_id)
            if not os.path.exists(fp):
                fp = os.path.join(os.environ.get("APPDATA",""), "open-webui","data","uploads", os.path.basename(fp))
            if not os.path.exists(fp):
                conn2.close()
                return json.dumps({"error": "File not found on disk"})
            conn2.close()
            src = fitz.open(fp)
            total_pages = src.page_count
            urls = []
            for start in range(0, total_pages, pages_per_file):
                end = min(start + pages_per_file, total_pages)
                sub = fitz.open()
                sub.insert_pdf(src, from_page=start, to_page=end-1)
                out = io.BytesIO()
                sub.save(out)
                sub.close()
                out.seek(0)
                part_name = f"part_{start+1}_{end}.pdf"
                url, name = self._save_and_link(out.read(), part_name)
                if url:
                    urls.append(f"[{name}]({url})")
            src.close()
            if urls:
                return "Split into " + str(len(urls)) + " files:\n" + "\n".join(urls)
            return json.dumps({"error": "Could not split PDF"})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    async def tool_stats(self, __user__=None, __request__=None) -> str:
        try:
            import sqlite3 as s3
            conn2 = s3.connect(r"C:\\Users\\Administrator\\AppData\\Roaming\\open-webui\\data\\webui.db")
            tool_count = conn2.execute("SELECT COUNT(*) FROM tool WHERE is_active=1").fetchone()[0]
            func_count = conn2.execute("SELECT COUNT(*) FROM function WHERE is_active=1").fetchone()[0]
            model_count = conn2.execute("SELECT COUNT(*) FROM model WHERE is_active=1").fetchone()[0]
            exports_dir = os.path.join(os.environ.get("LOCALAPPDATA", os.path.expanduser("~")), "open-webui", "exports")
            export_count = len([f for f in os.listdir(exports_dir) if os.path.isfile(os.path.join(exports_dir, f))]) if os.path.exists(exports_dir) else 0
            db_size_kb = os.path.getsize(r"C:\\Users\\Administrator\\AppData\\Roaming\\open-webui\\data\\webui.db") / 1024
            conn2.close()
            return json.dumps({
                "tools": tool_count,
                "functions": func_count,
                "models": model_count,
                "exported_files": export_count,
                "db_size_kb": round(db_size_kb, 1)
            }, indent=2)
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    
    async def manage_revisions(self, file_id: str, action: str, output_filename: str = "", __user__=None, __request__=None) -> str:
        """List, accept_all or reject_all tracked changes in a Word document."""
        try:
            import sqlite3 as s3
            conn2 = s3.connect(r"C:\\Users\\Administrator\\AppData\\Roaming\\open-webui\\data\\webui.db")
            row = conn2.execute("SELECT filename, meta FROM file WHERE id=?", (file_id,)).fetchone()
            if not row:
                row = conn2.execute("SELECT filename, meta FROM file WHERE filename LIKE ?", (f"%{file_id}%",)).fetchone()
            if not row:
                conn2.close()
                return json.dumps({"error": "File not found"})
            filename = row[0]
            meta = json.loads(row[1]) if row[1] else {}
            fp = meta.get("path", file_id)
            if not os.path.exists(fp):
                fp = os.path.join(os.environ.get("APPDATA",""), "open-webui", "data", "uploads", os.path.basename(fp))
            with open(fp, "rb") as f:
                data = f.read()
            conn2.close()
    
            from docx_revisions import RevisionDocument
    
            if action == "list":
                rdoc = RevisionDocument(io.BytesIO(data))
                revs = []
                for para in rdoc.paragraphs:
                    try:
                        rp = RevisionParagraph.from_paragraph(para)
                        if rp.has_track_changes:
                            for ins in rp.insertions:
                                revs.append({"type": "insertion", "author": ins.author, "text": ins.text[:100]})
                            for d in rp.deletions:
                                revs.append({"type": "deletion", "author": d.author, "text": d.text[:100]})
                    except Exception:
                        pass
                return json.dumps({"revisions": revs, "count": len(revs)}, indent=2)
    
            out_name = output_filename or filename
            rdoc = RevisionDocument(io.BytesIO(data))
            if action == "accept_all":
                rdoc.accept_all()
                msg = "All track changes accepted"
            elif action == "reject_all":
                rdoc.reject_all()
                msg = "All track changes rejected"
            else:
                return json.dumps({"error": f"Unknown action: {action}"})
    
            out = io.BytesIO()
            rdoc.save(out)
            out.seek(0)
            url, name = self._save_and_link(out.read(), out_name)
            if url:
                return f"[{name}]({url})\n\n{msg}."
            return json.dumps({"error": "Could not save file"})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

